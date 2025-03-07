import openai
from pptx import Presentation
import os
openai.api_key = "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"

def generate_slide_content(topic):
    prompt = f"Create a professional presentation outline on the topic: {topic}. Include slide titles and key points."
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "system", "content": "You are an expert presentation creator."},
                  {"role": "user", "content": prompt}]
    )
    return response["choices"][0]["message"]["content"]

def create_presentation(topic, filename="presentation.pptx"):
    content = generate_slide_content(topic)
    slides = content.split("\n\n")
    
    prs = Presentation()
    
    for slide in slides:
        lines = slide.split("\n")
        if not lines:
            continue
        
        slide_layout = prs.slide_layouts[1] 
        slide_obj = prs.slides.add_slide(slide_layout)
        title = slide_obj.shapes.title
        content = slide_obj.shapes.placeholders[1]
        
        title.text = lines[0]
        bullet_points = "\n".join(lines[1:])
        content.text = bullet_points
    
    prs.save(filename)
    print(f"Presentation saved as {filename}")

if __name__ == "__main__":
    topic = input("Enter the topic for the presentation: ")
    create_presentation(topic)